"""Document text extraction (PDF / DOCX / PPTX) with optional OCR.

PDF tables are extracted via `page.find_tables()` and serialized inline as
Markdown between `[TABLE N]…[/TABLE]` markers. Page images are passed through
Tesseract when available and emitted as `[IMAGE N OCR]…[/IMAGE]`.
"""
from pathlib import Path

from wisemock.config import (
    PDF_EXTRACT, DOCX_EXTRACT, PPTX_EXTRACT, OCR_EXTRACT, OCR_INSTALL_HINT,
)


# Tuning knobs for OCR fallback. Pages with little native text and large image
# coverage are likely scans; pure text pages skip OCR even if they contain
# decorative images.
OCR_FALLBACK_TEXT_THRESHOLD = 200
OCR_STRONG_TEXT_THRESHOLD = 500
SCANNED_IMAGE_COVERAGE_THRESHOLD = 0.55
OCR_RENDER_ZOOM = 2.0

# These optional libs are imported eagerly here only if their capability flag
# is True. Importing inside the function bodies would re-pay the cost on every
# extraction call.
if PDF_EXTRACT:
    import fitz as pymupdf
else:
    pymupdf = None

if DOCX_EXTRACT:
    import docx as python_docx
else:
    python_docx = None

if PPTX_EXTRACT:
    from pptx import Presentation as PptxPresentation
else:
    PptxPresentation = None

if OCR_EXTRACT:
    import io as _ocr_io
    import pytesseract
    from PIL import Image as _OcrImage
else:
    _ocr_io = None
    pytesseract = None
    _OcrImage = None


def _table_to_markdown(table) -> str:
    try:
        rows = table.extract()
    except Exception:
        return ""
    if not rows or not rows[0]:
        return ""
    header = [str(c or "").strip().replace("\n", " ") for c in rows[0]]
    out = ["| " + " | ".join(header) + " |",
           "| " + " | ".join("---" for _ in header) + " |"]
    for r in rows[1:]:
        cells = [str(c or "").strip().replace("\n", " ") for c in r]
        if len(cells) < len(header):
            cells += [""] * (len(header) - len(cells))
        else:
            cells = cells[:len(header)]
        out.append("| " + " | ".join(cells) + " |")
    return "\n".join(out)


def _extract_page_text_smart(page) -> str:
    """Extract page text reconstructing word boundaries from glyph coordinates.

    PyMuPDF's default `page.get_text()` can produce concatenated runs like
    "EUROPELARGESTMARITIMEPOWER" on slides where CSS-style letter-spacing
    spreads glyphs out: PDF stores each letter at its own position with no
    explicit space character, and the default extractor sometimes drops the
    inferred spaces. We fall back to word-level extraction (which carries
    coordinates) and, when a line looks letter-spaced (most "words" are
    single characters), cluster the glyphs by gap: small gaps stay intra-word,
    large gaps become spaces.
    """
    try:
        words_raw = page.get_text("words")
    except Exception:
        return page.get_text()
    if not words_raw:
        return page.get_text()
    # words tuple: (x0, y0, x1, y1, text, block_no, line_no, word_no)
    lines = {}
    line_order = []
    for w in words_raw:
        key = (w[5], w[6])
        if key not in lines:
            lines[key] = []
            line_order.append(key)
        lines[key].append(w)
    out_lines = []
    for key in line_order:
        items = sorted(lines[key], key=lambda w: w[0])
        if not items:
            continue
        single_char_ratio = sum(1 for w in items if len(w[4]) == 1) / len(items)
        if single_char_ratio > 0.5 and len(items) >= 4:
            # Letter-spaced line: cluster characters by inter-glyph gap.
            gaps = [items[i + 1][0] - items[i][2] for i in range(len(items) - 1)]
            sorted_gaps = sorted(gaps)
            median_gap = sorted_gaps[len(sorted_gaps) // 2] if sorted_gaps else 0
            # Anything noticeably wider than the median gap is a word boundary.
            threshold = max(median_gap * 1.8, 0.1)
            parts = [items[0][4]]
            for i, gap in enumerate(gaps):
                if gap > threshold:
                    parts.append(" ")
                parts.append(items[i + 1][4])
            out_lines.append("".join(parts))
        else:
            out_lines.append(" ".join(w[4] for w in items))
    return "\n".join(out_lines)


def _extraction_warning(message: str) -> str:
    return f"\n[EXTRACTION WARNING]\n{message}\n[/EXTRACTION WARNING]"


def _page_image_coverage(page) -> float:
    """Approximate how much of a PDF page is covered by raster images."""
    try:
        page_area = float(abs(page.rect.width * page.rect.height))
    except Exception:
        return 0.0
    if page_area <= 0:
        return 0.0
    covered = 0.0
    try:
        images = page.get_images(full=True)
    except Exception:
        return 0.0
    seen_rects = set()
    for img_info in images:
        xref = img_info[0]
        try:
            rects = page.get_image_rects(xref)
        except Exception:
            rects = []
        for rect in rects:
            key = (round(rect.x0, 1), round(rect.y0, 1), round(rect.x1, 1), round(rect.y1, 1))
            if key in seen_rects:
                continue
            seen_rects.add(key)
            covered += max(0.0, float(abs(rect.width * rect.height)))
    return min(1.0, covered / page_area)


def _should_ocr_page(page, page_text: str) -> tuple[bool, float]:
    """Return whether OCR is worth trying, plus image coverage for warnings."""
    text_len = len((page_text or "").strip())
    coverage = _page_image_coverage(page)
    if text_len >= OCR_STRONG_TEXT_THRESHOLD:
        return False, coverage
    if coverage >= SCANNED_IMAGE_COVERAGE_THRESHOLD and text_len < OCR_FALLBACK_TEXT_THRESHOLD:
        return True, coverage
    # Some scans are stored as one image with no reliable rect metadata.
    try:
        has_images = bool(page.get_images(full=True))
    except Exception:
        has_images = False
    return bool(has_images and text_len == 0), coverage


def _ocr_rendered_page(page) -> str:
    if not OCR_EXTRACT:
        return ""
    try:
        matrix = pymupdf.Matrix(OCR_RENDER_ZOOM, OCR_RENDER_ZOOM)
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        img = _OcrImage.open(_ocr_io.BytesIO(pix.tobytes("png")))
        return pytesseract.image_to_string(img).strip()
    except Exception:
        return ""


def _ocr_page_images(page, doc) -> list:
    if not OCR_EXTRACT:
        return []
    texts = []
    try:
        images = page.get_images(full=True)
    except Exception:
        return []
    for img_info in images:
        xref = img_info[0]
        try:
            pix = pymupdf.Pixmap(doc, xref)
            if pix.n - pix.alpha >= 4:
                pix = pymupdf.Pixmap(pymupdf.csRGB, pix)
            png_bytes = pix.tobytes("png")
            img = _OcrImage.open(_ocr_io.BytesIO(png_bytes))
            text = pytesseract.image_to_string(img).strip()
            if len(text) > 20:
                texts.append(text)
        except Exception:
            continue
    return texts


def _notify_progress(progress, percent, message):
    if progress:
        progress(max(0, min(100, int(percent))), message)


def extract_text_from_file(file_path: str, progress=None) -> str:
    ext = Path(file_path).suffix.lower()
    name = Path(file_path).name
    if ext == ".pdf":
        if not PDF_EXTRACT:
            raise ImportError("Install PyMuPDF to read PDFs:  pip install pymupdf")
        _notify_progress(progress, 1, f"Opening {name}...")
        doc = pymupdf.open(file_path)
        pages = []
        total_pages = max(getattr(doc, "page_count", 0), 1)
        for index, page in enumerate(doc, 1):
            _notify_progress(progress, ((index - 1) / total_pages) * 100,
                             f"Reading page {index} of {total_pages}...")
            page_text = _extract_page_text_smart(page)
            blocks = [page_text]
            try:
                tables = page.find_tables()
                table_list = list(tables) if tables else []
            except Exception:
                table_list = []
            for i, t in enumerate(table_list, 1):
                md = _table_to_markdown(t)
                if md:
                    blocks.append(f"\n[TABLE {i}]\n{md}\n[/TABLE]")
            # OCR is expensive (2-5s per image). Only run it on pages where
            # native text extraction yielded almost nothing and the page looks
            # image-backed. Text-heavy PDFs skip OCR even if they have logos.
            should_ocr, image_coverage = _should_ocr_page(page, page_text)
            if should_ocr:
                _notify_progress(progress, ((index - 0.35) / total_pages) * 100,
                                 f"Running OCR on page {index} of {total_pages}...")
                if OCR_EXTRACT:
                    rendered_text = _ocr_rendered_page(page)
                    if len(rendered_text) > 20:
                        blocks.append(f"\n[IMAGE {index} OCR]\n{rendered_text}\n[/IMAGE]")
                    elif not page_text.strip():
                        for i, ocr_txt in enumerate(_ocr_page_images(page, doc), 1):
                            blocks.append(f"\n[IMAGE {i} OCR]\n{ocr_txt}\n[/IMAGE]")
                else:
                    blocks.append(
                        _extraction_warning(
                            f"Page {index} appears to be scanned or image-based "
                            f"(image coverage ~{image_coverage:.0%}), but OCR is not available. "
                            f"{OCR_INSTALL_HINT}"
                        )
                    )
            pages.append("\n".join(blocks))
            _notify_progress(progress, (index / total_pages) * 100,
                             f"Finished page {index} of {total_pages}.")
        doc.close()
        _notify_progress(progress, 100, f"Finished extracting {name}.")
        return "\n\n".join(pages)
    elif ext == ".docx":
        if not DOCX_EXTRACT:
            raise ImportError("Install python-docx to read DOCX:  pip install python-docx")
        _notify_progress(progress, 15, f"Reading {name}...")
        doc = python_docx.Document(file_path)
        _notify_progress(progress, 100, f"Finished extracting {name}.")
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    elif ext == ".pptx":
        if not PPTX_EXTRACT:
            raise ImportError("Install python-pptx to read PPTX:  pip install python-pptx")
        _notify_progress(progress, 10, f"Reading {name}...")
        prs = PptxPresentation(file_path)
        slides_text = []
        total_slides = max(len(prs.slides), 1)
        for i, slide in enumerate(prs.slides, 1):
            _notify_progress(progress, ((i - 1) / total_slides) * 100,
                             f"Reading slide {i} of {total_slides}...")
            parts = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
            slides_text.append("\n".join(parts))
        _notify_progress(progress, 100, f"Finished extracting {name}.")
        return "\n\n".join(slides_text)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def combine_source_texts(sources: list[tuple[str, str]]) -> str:
    """Combine extracted study documents with stable provenance headers."""
    parts = []
    for index, (path, text) in enumerate(sources, 1):
        name = Path(path).name
        body = (text or "").strip()
        parts.append(f"--- SOURCE FILE {index}: {name} ---\n{body}")
    return "\n\n".join(parts).strip()
